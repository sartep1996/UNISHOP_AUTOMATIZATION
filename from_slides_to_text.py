from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    all_text = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

        all_text.append(f"Slide {slide_number}:\n" + "\n".join(text))

    return "\n\n".join(all_text)

pptx_path = "your_presentation.pptx"
extracted_text = extract_text_from_pptx(pptx_path)

with open("output.txt", "w", encoding="utf-8") as file:
    file.write(extracted_text)
